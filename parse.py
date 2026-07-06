"""文件解析工具 —— Python 实现，比 JS 库稳定"""
import sys
import json
import os
import zipfile
import re

def parse_txt(filepath):
    with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()

def parse_docx(filepath):
    try:
        from docx import Document
        doc = Document(filepath)
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    except ImportError:
        pass
    # fallback: 直接读 zip 中的 XML
    try:
        with zipfile.ZipFile(filepath) as z:
            xml = z.read("word/document.xml").decode("utf-8", errors="ignore")
            texts = re.findall(r"<w:t[^>]*>([^<]*)</w:t>", xml)
            return "".join(texts)
    except:
        raise Exception("无法解析 DOCX 文件，请确认文件未被损坏")

def parse_pdf(filepath):
    try:
        import fitz  # PyMuPDF
        doc = fitz.open(filepath)
        text = "\n\n".join(page.get_text() for page in doc)
        doc.close()
        return text
    except ImportError:
        pass
    try:
        from PyPDF2 import PdfReader
        reader = PdfReader(filepath)
        return "\n\n".join(page.extract_text() or "" for page in reader.pages)
    except ImportError:
        raise Exception("缺少 PDF 解析库，请安装: pip install PyMuPDF")

def parse_pptx(filepath):
    try:
        from pptx import Presentation
        prs = Presentation(filepath)
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    texts.append(shape.text_frame.text)
        return "\n".join(texts)
    except ImportError:
        pass
    # fallback: zip + xml
    try:
        with zipfile.ZipFile(filepath) as z:
            slides = [f for f in z.namelist() if f.startswith("ppt/slides/slide") and f.endswith(".xml")]
            texts = []
            for s in sorted(slides):
                xml = z.read(s).decode("utf-8", errors="ignore")
                found = re.findall(r"<a:t[^>]*>([^<]*)</a:t>", xml)
                texts.append("".join(found))
            return "\n".join(texts)
    except:
        raise Exception("无法解析 PPTX 文件")

def parse_zip(filepath):
    with zipfile.ZipFile(filepath) as z:
        all_text = []
        for name in z.namelist():
            ext = os.path.splitext(name)[1].lower()
            if ext not in (".txt", ".md", ".docx", ".pdf", ".pptx"):
                continue
            try:
                data = z.read(name)
                tmp = os.path.join(os.path.dirname(filepath), f"_tmp_{os.path.basename(name)}")
                with open(tmp, "wb") as f:
                    f.write(data)
                result = parse_file(tmp, os.path.splitext(name)[1].lstrip("."))
                os.unlink(tmp)
                all_text.append(f"--- {name} ---\n{result}")
            except:
                pass
        return "\n\n".join(all_text)

def parse_doc(filepath):
    """解析旧版 .doc 文件（用 antiword 转换）"""
    import subprocess
    import shutil
    antiword = shutil.which("antiword")
    if not antiword:
        raise Exception("不支持旧版 .doc 格式，请用 Word / WPS 另存为 .docx 后重新上传")
    try:
        result = subprocess.run([antiword, "-m", "UTF-8.txt", filepath],
                                capture_output=True, text=True, timeout=15, encoding="utf-8", errors="ignore")
        text = result.stdout.strip()
        if not text or len(text) < 10:
            raise Exception("antiword 提取内容为空")
        return text
    except Exception as e:
        raise Exception(f"旧版 .doc 解析失败: {e}")

def parse_file(filepath, ext):
    if ext in ("txt", "md"):
        return parse_txt(filepath)
    if ext == "pdf":
        return parse_pdf(filepath)
    if ext == "docx":
        return parse_docx(filepath)
    if ext == "doc":
        return parse_doc(filepath)
    if ext == "pptx":
        return parse_pptx(filepath)
    if ext == "zip":
        return parse_zip(filepath)
    raise Exception(f"不支持的格式: .{ext}")

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print(json.dumps({"error": "缺少文件路径参数"}))
        sys.exit(1)
    filepath = sys.argv[1]
    ext = os.path.splitext(filepath)[1].lower().lstrip(".")
    try:
        text = parse_file(filepath, ext)
        if not text or len(text.strip()) < 10:
            print(json.dumps({"error": f"文件内容太少（{len(text)}字符），可能是扫描件"}))
        else:
            print(json.dumps({"text": text.strip()}))
    except Exception as e:
        print(json.dumps({"error": str(e)}))
        sys.exit(1)
