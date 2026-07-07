"""Vercel Python Runtime: /api/parse 文件解析端点"""
import json, os, tempfile, re, zipfile, subprocess, shutil
from http.server import BaseHTTPRequestHandler
from io import BytesIO

# ======== 解析函数（内联，不依赖外部 parse.py）========

def parse_txt(filepath):
    with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()

def parse_docx(filepath):
    try:
        from docx import Document
        doc = Document(filepath)
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    except:
        with zipfile.ZipFile(filepath) as z:
            xml = z.read("word/document.xml").decode("utf-8", errors="ignore")
            return "".join(re.findall(r"<w:t[^>]*>([^<]*)</w:t>", xml))

def parse_pdf(filepath):
    try:
        import fitz
        doc = fitz.open(filepath)
        text = "\n\n".join(page.get_text() for page in doc)
        doc.close()
        return text
    except:
        from PyPDF2 import PdfReader
        reader = PdfReader(filepath)
        return "\n\n".join(page.extract_text() or "" for page in reader.pages)

def parse_pptx(filepath):
    try:
        from pptx import Presentation
        prs = Presentation(filepath)
        return "\n".join(
            shape.text for slide in prs.slides
            for shape in slide.shapes if shape.has_text_frame
        )
    except:
        with zipfile.ZipFile(filepath) as z:
            texts = []
            for s in sorted(f for f in z.namelist() if f.startswith("ppt/slides/slide") and f.endswith(".xml")):
                xml = z.read(s).decode("utf-8", errors="ignore")
                texts.append("".join(re.findall(r"<a:t[^>]*>([^<]*)</a:t>", xml)))
            return "\n".join(texts)

def parse_doc(filepath):
    antiword = shutil.which("antiword")
    if antiword:
        r = subprocess.run([antiword, filepath], capture_output=True, text=True, timeout=15)
        text = r.stdout.strip()
        if text:
            return text
    raise Exception("旧版 .doc 请用 Word 另存为 .docx 后上传")

def parse_file(filepath, ext):
    if ext in ("txt", "md"): return parse_txt(filepath)
    if ext == "pdf": return parse_pdf(filepath)
    if ext == "docx": return parse_docx(filepath)
    if ext == "doc": return parse_doc(filepath)
    if ext == "pptx": return parse_pptx(filepath)
    raise Exception(f"不支持的格式: .{ext}")

# ======== Vercel Handler ========

class handler(BaseHTTPRequestHandler):
    def do_POST(self):
        try:
            content_type = self.headers.get("Content-Type", "")
            if "multipart/form-data" not in content_type:
                self._json({"error": "需要 multipart/form-data"}, 400)
                return

            # 解析 multipart
            boundary = content_type.split("boundary=")[1]
            body = self.rfile.read(int(self.headers.get("Content-Length", 0)))

            # 找 file 字段
            file_data, filename = self._extract_file(body, boundary.encode())
            if not file_data or not filename:
                self._json({"error": "没有接收到文件"}, 400)
                return

            ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
            if ext not in ("txt", "md", "pdf", "docx", "doc", "pptx"):
                self._json({"error": f"不支持的格式: .{ext}，支持 txt/md/pdf/docx/doc/pptx"}, 400)
                return

            with tempfile.NamedTemporaryFile(suffix=f".{ext}", delete=False) as tmp:
                tmp.write(file_data)
                tmp_path = tmp.name

            try:
                text = parse_file(tmp_path, ext)
                if not text or len(text.strip()) < 10:
                    self._json({"error": "文件内容太少，可能是扫描件或受保护的文档"}, 422)
                else:
                    self._json({"text": text.strip(), "filename": filename})
            finally:
                try: os.unlink(tmp_path)
                except: pass

        except Exception as e:
            self._json({"error": f"解析失败：{e}"}, 500)

    def _extract_file(self, body, boundary):
        """从 multipart body 中提取 file 字段的内容"""
        parts = body.split(b"--" + boundary)
        for part in parts:
            if b'name="file"' in part:
                # 找到 Content-Type 后的空行
                header_end = part.find(b"\r\n\r\n")
                if header_end == -1:
                    header_end = part.find(b"\n\n")
                if header_end == -1:
                    continue
                raw = part[header_end:].strip()
                # 去掉末尾的 boundary
                end = raw.rfind(b"\r\n--")
                if end == -1:
                    end = raw.rfind(b"\n--")
                if end > 0:
                    raw = raw[:end]
                # 提取文件名
                fn_match = re.search(rb'filename="([^"]*)"', part)
                filename = fn_match.group(1).decode("utf-8", errors="ignore") if fn_match else "unknown"
                return raw.strip(), filename
        return None, None

    def _json(self, data, status=200):
        body = json.dumps(data, ensure_ascii=False).encode("utf-8")
        self.send_response(status)
        self.send_header("Content-Type", "application/json; charset=utf-8")
        self.send_header("Content-Length", str(len(body)))
        self.end_headers()
        self.wfile.write(body)
